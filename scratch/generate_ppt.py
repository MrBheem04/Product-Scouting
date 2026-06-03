import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Initialize presentation
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
DARK_BG = RGBColor(15, 23, 42)      # #0F172A
LIGHT_TEXT = RGBColor(226, 232, 240) # #E2E8F0
PURPLE_ACCENT = RGBColor(168, 85, 247) # #A855F7
BLUE_ACCENT = RGBColor(59, 130, 246)   # #3B82F6
MUTED_TEXT = RGBColor(148, 163, 184)   # #94A6B8

# Image paths (Absolute paths from the brain directory)
brain_dir = r"C:\Users\LENOVO\.gemini\antigravity\brain\842f72d0-de2a-42b7-87bf-b090b6b83e64"
hero_img = os.path.join(brain_dir, "scoutprice_hero_1780110489127.png")
ext_img = os.path.join(brain_dir, "scoutprice_extension_ui_1780110531680.png")
chat_img = os.path.join(brain_dir, "scoutprice_chatbot_ui_1780110550935.png")

def set_slide_background(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG

def add_title(slide, text, color=PURPLE_ACCENT):
    # Add title text box
    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.833), Inches(1.0))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = color
    p.font.name = "Arial"
    return title_box

# ==========================================
# SLIDE 1: Title Slide (Blank Layout 6)
# ==========================================
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)
set_slide_background(slide)

# Main Title
title_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.333), Inches(2.0))
tf = title_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "ScoutPrice 🚀"
p.font.size = Pt(64)
p.font.bold = True
p.font.color.rgb = PURPLE_ACCENT
p.alignment = PP_ALIGN.CENTER

p2 = tf.add_paragraph()
p2.text = "AI-Powered Price Tracker & Shopping Assistant"
p2.font.size = Pt(28)
p2.font.bold = True
p2.font.color.rgb = BLUE_ACCENT
p2.alignment = PP_ALIGN.CENTER

# Subtitle Info
info_box = slide.shapes.add_textbox(Inches(1.0), Inches(4.8), Inches(11.333), Inches(1.5))
tf_info = info_box.text_frame
p3 = tf_info.paragraphs[0]
p3.text = "Next-Gen Price Monitoring, Referral Shielding, & AI Shopping Companion"
p3.font.size = Pt(16)
p3.font.color.rgb = LIGHT_TEXT
p3.alignment = PP_ALIGN.CENTER

p4 = tf_info.add_paragraph()
p4.text = "Stack: React • Node.js • Chrome MV3 • Redis • MongoDB • Nginx"
p4.font.size = Pt(14)
p4.font.color.rgb = MUTED_TEXT
p4.alignment = PP_ALIGN.CENTER


# ==========================================
# SLIDE 2: The E-Commerce Chaos (Problem)
# ==========================================
slide = prs.slides.add_slide(slide_layout)
set_slide_background(slide)
add_title(slide, "The E-Commerce Chaos")

# Left Column - Issues list
content_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.833), Inches(5.0))
tf = content_box.text_frame
tf.word_wrap = True

points = [
    ("Dynamic Pricing Manipulation", "Retailers continuously adjust prices in real-time, exploiting consumer urgency and browser cookies."),
    ("Store Fragmentation", "Comparing prices across Amazon, Flipkart, Myntra, Ajio, and Meesho is manual, slow, and frustrating."),
    ("Fake Coupons & Stale Offers", "Shoppers waste time testing invalid coupon codes manually during checkout."),
    ("Anti-Scraping Barriers", "Platforms block automated checkers with security walls, leading to error landing screens (e.g. Flipkart E002).")
]

for title, desc in points:
    p = tf.add_paragraph()
    p.text = f"•  {title}"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = BLUE_ACCENT
    p.space_after = Pt(4)
    
    p_desc = tf.add_paragraph()
    p_desc.text = f"    {desc}"
    p_desc.font.size = Pt(16)
    p_desc.font.color.rgb = LIGHT_TEXT
    p_desc.space_after = Pt(20)


# ==========================================
# SLIDE 3: The ScoutPrice Solution
# ==========================================
slide = prs.slides.add_slide(slide_layout)
set_slide_background(slide)
add_title(slide, "The ScoutPrice Solution")

content_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.833), Inches(5.0))
tf = content_box.text_frame
tf.word_wrap = True

solutions = [
    ("Unified Scraper Workers", "Dual-mode Puppeteer + Cheerio parser fetches accurate product data scheduled via Redis / BullMQ."),
    ("Glassmorphism Dashboard", "A sleek dark-themed analytics portal presenting history charts, lowest/highest price indicators, and buy recommendations."),
    ("Automated Coupon Application", "A Chrome Extension background script that loops through and auto-solves valid discount codes on the checkout page."),
    ("Referrer Shield Proxy", "Built-in server redirect gateway to strip referrers, bypass browser cache, and avoid store blockers.")
]

for title, desc in solutions:
    p = tf.add_paragraph()
    p.text = f"✔  {title}"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = PURPLE_ACCENT
    p.space_after = Pt(4)
    
    p_desc = tf.add_paragraph()
    p_desc.text = f"    {desc}"
    p_desc.font.size = Pt(16)
    p_desc.font.color.rgb = LIGHT_TEXT
    p_desc.space_after = Pt(20)


# ==========================================
# SLIDE 4: Architecture & Technical Stack
# ==========================================
slide = prs.slides.add_slide(slide_layout)
set_slide_background(slide)
add_title(slide, "Technical Stack & Architecture", color=BLUE_ACCENT)

# Left Box: Tech Stack
stack_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(5.5), Inches(4.8))
tf = stack_box.text_frame
tf.word_wrap = True

p_hdr = tf.paragraphs[0]
p_hdr.text = "Microservices Stack"
p_hdr.font.size = Pt(24)
p_hdr.font.bold = True
p_hdr.font.color.rgb = PURPLE_ACCENT
p_hdr.space_after = Pt(14)

stack_items = [
    ("Vite React + Redux Toolkit", "Renders the responsive state-synchronized dashboard UI."),
    ("Node.js + Express.js", "Core REST API server managing user sessions, watchlists, and price alerts."),
    ("Nginx Reverse Proxy", "Acts as the single port 80 entry gateway routing frontend & API services."),
    ("Redis + BullMQ", "Coordinates async scraper scheduling to prevent IP bans."),
    ("MongoDB (Mongoose)", "Houses product metadata, user alerts, and composite price history.")
]

for title, desc in stack_items:
    p = tf.add_paragraph()
    p.text = f"▪ {title}"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = LIGHT_TEXT
    p.space_after = Pt(2)
    
    p_desc = tf.add_paragraph()
    p_desc.text = f"  {desc}"
    p_desc.font.size = Pt(14)
    p_desc.font.color.rgb = MUTED_TEXT
    p_desc.space_after = Pt(10)

# Right Box: Layout description
right_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.8), Inches(5.5), Inches(4.8))
tf_right = right_box.text_frame
tf_right.word_wrap = True

p_rhdr = tf_right.paragraphs[0]
p_rhdr.text = "Key Infrastructure Flows"
p_rhdr.font.size = Pt(24)
p_rhdr.font.bold = True
p_rhdr.font.color.rgb = BLUE_ACCENT
p_rhdr.space_after = Pt(14)

flows = [
    ("Chrome Extension MV3 Helper", "Injects active content script widgets to extract ASINs, titles, and match cross-store prices dynamically."),
    ("WebSockets (Socket.io)", "Pushes instant, low-latency price-drop alerts directly to connected user browser viewports."),
    ("Anti-Bot Redirection Shield", "Strips HTTP referrer tokens during merchant handoffs to prevent crawler block screens.")
]

for title, desc in flows:
    p = tf_right.add_paragraph()
    p.text = f"▸ {title}"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = LIGHT_TEXT
    p.space_after = Pt(2)
    
    p_desc = tf_right.add_paragraph()
    p_desc.text = f"  {desc}"
    p_desc.font.size = Pt(14)
    p_desc.font.color.rgb = MUTED_TEXT
    p_desc.space_after = Pt(10)


# ==========================================
# SLIDE 5: Web Dashboard & Price Tracker
# ==========================================
slide = prs.slides.add_slide(slide_layout)
set_slide_background(slide)
add_title(slide, "Sleek Dashboard & Analytics")

# Details on Left
detail_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(5.5), Inches(4.8))
tf = detail_box.text_frame
tf.word_wrap = True

bullets = [
    ("Visual Analytics", "Clean charts powered by Recharts plotting detailed price variation graphs over time."),
    ("Price Level Gauges", "Highlights minimum, maximum, and average tags relative to the active price."),
    ("Buy Recommendations", "Calculates a target score determining whether to buy immediately or wait for price drops."),
    ("Integrated Watchlist", "Allows shoppers to save products and specify discount target thresholds.")
]

for title, desc in bullets:
    p = tf.add_paragraph() if tf.paragraphs[0].text else tf.paragraphs[0]
    p.text = f"• {title}"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = BLUE_ACCENT
    p.space_after = Pt(2)
    
    p_desc = tf.add_paragraph()
    p_desc.text = f"  {desc}"
    p_desc.font.size = Pt(15)
    p_desc.font.color.rgb = LIGHT_TEXT
    p_desc.space_after = Pt(14)

# Image on Right
if os.path.exists(hero_img):
    slide.shapes.add_picture(hero_img, Inches(6.5), Inches(1.8), width=Inches(6.2))


# ==========================================
# SLIDE 6: Smart Redirection Proxy
# ==========================================
slide = prs.slides.add_slide(slide_layout)
set_slide_background(slide)
add_title(slide, "Referrer Shield & Redirection", color=BLUE_ACCENT)

content_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.833), Inches(5.0))
tf = content_box.text_frame
tf.word_wrap = True

details = [
    ("The Redirect Challenge", "E-commerce gateways (like Flipkart) flag requests initiated from 'localhost' development origins, showing bot block screens."),
    ("Referrer Stripping Engine", "By running outbound links through `/api/affiliate/redirect` and setting the response header `Referrer-Policy: no-referrer`, we strip trace vectors."),
    ("Strict Browser Cache Bypass", "Ensures browsers do not cache failed redirects by implementing anti-caching headers (Cache-Control: no-store, no-cache)."),
    ("Dynamic SKU Resolvers", "Cross-platform token-based matching checks catalog entries to serve direct target pages (avoiding generic queries).")
]

for title, desc in details:
    p = tf.add_paragraph() if tf.paragraphs[0].text else tf.paragraphs[0]
    p.text = f"▪ {title}"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = PURPLE_ACCENT
    p.space_after = Pt(2)
    
    p_desc = tf.add_paragraph()
    p_desc.text = f"  {desc}"
    p_desc.font.size = Pt(15)
    p_desc.font.color.rgb = LIGHT_TEXT
    p_desc.space_after = Pt(14)


# ==========================================
# SLIDE 7: Chrome Extension (Manifest V3)
# ==========================================
slide = prs.slides.add_slide(slide_layout)
set_slide_background(slide)
add_title(slide, "Chrome Extension Helper")

# Details on Left
detail_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(5.5), Inches(4.8))
tf = detail_box.text_frame
tf.word_wrap = True

bullets_ext = [
    ("Zero-Click Automation", "Runs content scripts in the background, executing code only when navigating matching store portals."),
    ("On-Page Comparison Card", "Injects a floating sidebar comparing current store pricing against competing merchants instantly."),
    ("Coupon Autofill Engine", "Gathers verified database codes and executes auto-type validation routines on the checkout page."),
    ("User Preference Control", "Minimal footprint widget that can be minimized or toggled without page reloads.")
]

for title, desc in bullets_ext:
    p = tf.add_paragraph() if tf.paragraphs[0].text else tf.paragraphs[0]
    p.text = f"• {title}"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = PURPLE_ACCENT
    p.space_after = Pt(2)
    
    p_desc = tf.add_paragraph()
    p_desc.text = f"  {desc}"
    p_desc.font.size = Pt(15)
    p_desc.font.color.rgb = LIGHT_TEXT
    p_desc.space_after = Pt(14)

# Image on Right
if os.path.exists(ext_img):
    slide.shapes.add_picture(ext_img, Inches(6.5), Inches(1.8), width=Inches(6.2))


# ==========================================
# SLIDE 8: AI Chatbot Companion
# ==========================================
slide = prs.slides.add_slide(slide_layout)
set_slide_background(slide)
add_title(slide, "AI Shopping Chatbot", color=BLUE_ACCENT)

# Details on Left
detail_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(5.5), Inches(4.8))
tf = detail_box.text_frame
tf.word_wrap = True

bullets_chat = [
    ("Contextual Recommendations", "Handles user questions about best timings, low prices, and specs via conversational interfaces."),
    ("Expert Fallback Engines", "Structured rule logic triggers if LLM API connections timeout, assuring 100% service availability."),
    ("Case-Insensitive Searching", "Corrects casing checks (e.g. 'electronics' to 'Electronics') dynamically to avoid empty collections queries."),
    ("Interactive Data Snippets", "Presents mini price logs and best store suggestion cards inside user chat bubbles.")
]

for title, desc in bullets_chat:
    p = tf.add_paragraph() if tf.paragraphs[0].text else tf.paragraphs[0]
    p.text = f"✔ {title}"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = BLUE_ACCENT
    p.space_after = Pt(2)
    
    p_desc = tf.add_paragraph()
    p_desc.text = f"  {desc}"
    p_desc.font.size = Pt(15)
    p_desc.font.color.rgb = LIGHT_TEXT
    p_desc.space_after = Pt(14)

# Image on Right
if os.path.exists(chat_img):
    slide.shapes.add_picture(chat_img, Inches(6.5), Inches(1.8), width=Inches(6.2))


# ==========================================
# SLIDE 9: Database Schema & Collections
# ==========================================
slide = prs.slides.add_slide(slide_layout)
set_slide_background(slide)
add_title(slide, "Mongoose Database Design")

# Left Column (Entities list)
schema_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.833), Inches(5.0))
tf = schema_box.text_frame
tf.word_wrap = True

schemas = [
    ("Users Collection", "Encrypted hash logins, saved settings, and active watchlist lists."),
    ("Products Collection", "SKU identifiers, brand info, and merchant details across Amazon, Flipkart, Myntra, Ajio, and Meesho."),
    ("PriceHistory Collection", "Plot coordinates mapping timestamps to actual product pricing for detailed dashboard charting."),
    ("Alerts Collection", "Tracks target trigger prices, alert channels (Email / WebSockets), and status flags."),
    ("Fitness & Sports Category", "Added category mapping and catalog additions (e.g. Kore DM 20kg Dumbbell Sets) synced successfully.")
]

for title, desc in schemas:
    p = tf.add_paragraph() if tf.paragraphs[0].text else tf.paragraphs[0]
    p.text = f"▪ {title}"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = BLUE_ACCENT
    p.space_after = Pt(2)
    
    p_desc = tf.add_paragraph()
    p_desc.text = f"  {desc}"
    p_desc.font.size = Pt(14)
    p_desc.font.color.rgb = LIGHT_TEXT
    p_desc.space_after = Pt(10)


# ==========================================
# SLIDE 10: Technical Accomplishments (Summary)
# ==========================================
slide = prs.slides.add_slide(slide_layout)
set_slide_background(slide)
add_title(slide, "Completed Technical Achievements", color=PURPLE_ACCENT)

content_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.833), Inches(5.0))
tf = content_box.text_frame
tf.word_wrap = True

achievements = [
    ("Affiliate Path Optimization", "Stripped referring headers and added cache controls to fix browser redirection and Flipkart E002 blocks."),
    ("Cross-Merchant Sync", "Configured real URLs and correct prices for products, filtering store listings based on category rules."),
    ("Database & Logic Alignment", "Fixed category case-sensitivity issue in search indexing and added the Fitness & Sports section."),
    ("UI Component Refinement", "Cleaned up charts layout spacing to prevent overlay elements colliding on product detail panels."),
    ("Integration and Testing", "Validated microservice setups via Docker and built a PowerShell REST API validation script.")
]

for title, desc in achievements:
    p = tf.add_paragraph() if tf.paragraphs[0].text else tf.paragraphs[0]
    p.text = f"✔ {title}"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = BLUE_ACCENT
    p.space_after = Pt(2)
    
    p_desc = tf.add_paragraph()
    p_desc.text = f"  {desc}"
    p_desc.font.size = Pt(14)
    p_desc.font.color.rgb = LIGHT_TEXT
    p_desc.space_after = Pt(10)

# Save presentation
output_path = r"c:\Users\LENOVO\OneDrive\Desktop\Product Scouting\ScoutPrice_Presentation.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
